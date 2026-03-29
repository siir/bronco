#!/usr/bin/env python3
"""
Generate the Bronco architecture presentation as a .pptx file.

Usage:
    python3 scripts/generate-architecture-pptx.py

Requirements:
    - python-pptx (e.g. ``pip install python-pptx==0.6.21``)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Colour palette ──────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x1E, 0x1E, 0x2E)   # dark navy background
BG_SECTION   = RGBColor(0x28, 0x28, 0x3C)   # slightly lighter section bg
ACCENT_BLUE  = RGBColor(0x58, 0x9C, 0xF5)   # primary accent
ACCENT_GREEN = RGBColor(0x6B, 0xCB, 0x77)   # success / green
ACCENT_ORANGE= RGBColor(0xF0, 0x9E, 0x4A)   # warning / orange
ACCENT_PURPLE= RGBColor(0xB3, 0x8B, 0xFA)   # purple accent
ACCENT_RED   = RGBColor(0xF0, 0x65, 0x65)   # red accent
ACCENT_TEAL  = RGBColor(0x4E, 0xC9, 0xB0)   # teal
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY     = RGBColor(0x99, 0x99, 0x99)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helper utilities ────────────────────────────────────────────────────

def add_bg(slide, color=BG_DARK):
    """Fill the entire slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(1)):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_title_subtitle(slide, title, subtitle=None, title_size=44):
    """Add a big title and optional subtitle."""
    add_text_box(slide, Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.2),
                 title, font_size=title_size, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.6),
                     subtitle, font_size=20, color=LIGHT_GRAY)

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=LIGHT_GRAY, spacing=Pt(6)):
    """Add a multi-line bullet list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txBox

def add_component_card(slide, left, top, width, height, title, items,
                       accent=ACCENT_BLUE, title_size=18, item_size=13):
    """Draw a labelled card with a coloured top border."""
    # Card background
    card = add_shape(slide, left, top, width, height, fill_color=BG_SECTION, border_color=RGBColor(0x40, 0x40, 0x55))
    # Accent stripe at top
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()
    stripe.shadow.inherit = False
    # Title
    add_text_box(slide, left + Inches(0.15), top + Inches(0.10), width - Inches(0.3), Inches(0.4),
                 title, font_size=title_size, color=accent, bold=True)
    # Items
    if items:
        add_bullet_list(slide, left + Inches(0.15), top + Inches(0.50), width - Inches(0.3),
                        height - Inches(0.60), items, font_size=item_size, color=LIGHT_GRAY, spacing=Pt(3))
    return card

def add_arrow_down(slide, left, top, height, color=ACCENT_BLUE):
    """Draw a small downward-pointing arrow."""
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, Inches(0.25), height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    arrow.shadow.inherit = False
    return arrow


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.5),
             "Bronco", font_size=60, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.8), Inches(3.3), Inches(11.7), Inches(1.0),
             "Architecture Overview", font_size=36, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.8), Inches(4.3), Inches(11.7), Inches(0.8),
             "AI-augmented database & software operations platform", font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.5),
             "TypeScript  \u00b7  pnpm monorepo  \u00b7  Docker Compose  \u00b7  PostgreSQL  \u00b7  BullMQ  \u00b7  Prisma",
             font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 2 — What is Bronco?
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_subtitle(slide, "What is Bronco?",
                   "Single-operator platform for managing client database systems & software operations")

items = [
    "\u2022  Manages client database systems (Azure SQL Managed Instances, on-prem SQL Server)",
    "\u2022  Triages tickets across: DB performance, bug fixes, features, code reviews, architecture",
    "\u2022  Ingests work from: Email (IMAP), Azure DevOps work items, manual API/UI, scheduled jobs",
    "\u2022  AI-powered analysis: Local Ollama for fast triage, Claude API for deep reasoning",
    "\u2022  Automated issue resolution: Claude generates code fixes and pushes to feature branches",
    "\u2022  YouTube scheduler: Automates weekly worship service live stream setup",
    "\u2022  Control panel: Angular SPA behind Caddy reverse proxy (Tailscale access)",
]
add_bullet_list(slide, Inches(0.8), Inches(2.6), Inches(11.7), Inches(4.5),
                items, font_size=18, color=LIGHT_GRAY, spacing=Pt(10))

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 3 — High-Level Architecture Diagram
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_subtitle(slide, "High-Level Architecture", "How the pieces fit together")

# --- External sources (top row) ---
ext_y = Inches(2.4)
ext_h = Inches(0.9)

add_component_card(slide, Inches(0.3), ext_y, Inches(2.0), ext_h,
                   "Email (IMAP)", ["\u2709 Gmail / IMAP server"], accent=ACCENT_ORANGE, title_size=14, item_size=11)
add_component_card(slide, Inches(2.6), ext_y, Inches(2.0), ext_h,
                   "Azure DevOps", ["Work items & comments"], accent=ACCENT_ORANGE, title_size=14, item_size=11)
add_component_card(slide, Inches(5.0), ext_y, Inches(2.0), ext_h,
                   "GitHub", ["Issues & repos"], accent=ACCENT_ORANGE, title_size=14, item_size=11)
add_component_card(slide, Inches(7.4), ext_y, Inches(2.0), ext_h,
                   "Google Drive", [".docx bulletins"], accent=ACCENT_ORANGE, title_size=14, item_size=11)
add_component_card(slide, Inches(9.8), ext_y, Inches(2.3), ext_h,
                   "Client SQL Servers", ["Azure SQL MI / on-prem"], accent=ACCENT_ORANGE, title_size=14, item_size=11)

# Arrows down from external sources
for x_pos in [Inches(1.1), Inches(3.4), Inches(5.8), Inches(8.2)]:
    add_arrow_down(slide, x_pos, Inches(3.35), Inches(0.35), ACCENT_ORANGE)

# --- Workers row ---
worker_y = Inches(3.8)
worker_h = Inches(1.2)

add_component_card(slide, Inches(0.3), worker_y, Inches(2.0), worker_h,
                   "imap-worker", ["Polls IMAP", "Creates tickets", "Threads emails"], accent=ACCENT_GREEN, title_size=14, item_size=11)
add_component_card(slide, Inches(2.6), worker_y, Inches(2.0), worker_h,
                   "devops-worker", ["Polls Azure DevOps", "AI conversations", "Plan execution"], accent=ACCENT_GREEN, title_size=14, item_size=11)
add_component_card(slide, Inches(5.0), worker_y, Inches(2.0), worker_h,
                   "issue-resolver", ["Clones repos", "Claude code gen", "Push to branch"], accent=ACCENT_GREEN, title_size=14, item_size=11)
add_component_card(slide, Inches(7.4), worker_y, Inches(2.0), worker_h,
                   "yt-scheduler", ["Bulletin \u2192 PDF", "Gen thumbnail", "Schedule stream"], accent=ACCENT_GREEN, title_size=14, item_size=11)

# MCP server (connects to client DBs)
add_component_card(slide, Inches(9.8), worker_y, Inches(2.3), worker_h,
                   "MCP DB Server", ["Azure App Service", "SQL tools (MCP)", "Pool manager"], accent=ACCENT_TEAL, title_size=14, item_size=11)
add_arrow_down(slide, Inches(10.75), Inches(3.35), Inches(0.35), ACCENT_TEAL)

# Arrows down to shared layer
for x_pos in [Inches(1.1), Inches(3.4), Inches(5.8)]:
    add_arrow_down(slide, x_pos, Inches(5.05), Inches(0.35), ACCENT_GREEN)

# --- Shared layer ---
shared_y = Inches(5.5)

add_component_card(slide, Inches(0.3), shared_y, Inches(3.3), Inches(1.5),
                   "copilot-api (Fastify)", ["REST API  \u00b7  Port 3000", "Ticket CRUD  \u00b7  BullMQ producer", "Auth (JWT + API key)"], accent=ACCENT_BLUE, title_size=14, item_size=11)

add_component_card(slide, Inches(3.9), shared_y, Inches(2.3), Inches(1.5),
                   "PostgreSQL", ["Control plane DB", "Prisma ORM", "Tickets, clients, events"], accent=ACCENT_PURPLE, title_size=14, item_size=11)

add_component_card(slide, Inches(6.5), shared_y, Inches(2.0), Inches(1.5),
                   "Redis", ["BullMQ job queues", "Worker coordination"], accent=ACCENT_RED, title_size=14, item_size=11)

add_component_card(slide, Inches(8.8), shared_y, Inches(2.0), Inches(1.5),
                   "Ollama (Mac mini)", ["Local LLM inference", "Triage, classify", "Summarize, draft"], accent=ACCENT_PURPLE, title_size=14, item_size=11)

add_component_card(slide, Inches(11.1), shared_y, Inches(1.9), Inches(1.5),
                   "Claude API", ["Deep analysis", "Code generation", "Architecture review"], accent=ACCENT_PURPLE, title_size=14, item_size=11)

# Connecting arrows: workers to Postgres/Redis
add_arrow_down(slide, Inches(4.9), Inches(5.05), Inches(0.35), ACCENT_PURPLE)

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Monorepo Structure
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_subtitle(slide, "Monorepo Structure", "pnpm workspaces  \u00b7  TypeScript throughout  \u00b7  ESM imports")

col1_x = Inches(0.5)
col2_x = Inches(4.6)
col3_x = Inches(9.0)
card_w = Inches(3.8)
top_y = Inches(2.5)

# packages/
add_component_card(slide, col1_x, top_y, card_w, Inches(4.2),
                   "packages/  (shared libraries)", [
                       "shared-types \u2500 enums, interfaces",
                       "    const object + type pattern (no TS enums)",
                       "    TicketStatus, Priority, TaskType, etc.",
                       "",
                       "shared-utils \u2500 logging, config, crypto",
                       "    createLogger(name) \u2192 Pino to stderr",
                       "    loadConfig(zodSchema) \u2192 env validation",
                       "    encrypt/decrypt (AES-256-GCM) for credentials",
                       "    AppLogger \u2192 structured logs to Postgres",
                       "",
                       "db \u2500 Prisma ORM + migrations",
                       "    schema.prisma defines all models",
                       "    ensureClientUser() auto-provisions users",
                       "",
                       "ai-provider \u2500 AI routing layer",
                       "    AIRouter: Ollama vs Claude routing",
                       "    OllamaClient / ClaudeClient adapters",
                   ], accent=ACCENT_PURPLE, title_size=16, item_size=12)

# services/
add_component_card(slide, col2_x, top_y, card_w, Inches(4.2),
                   "services/  (deployable apps)", [
                       "copilot-api \u2500 Fastify REST API",
                       "    Ticket, client, system CRUD",
                       "    BullMQ job producer",
                       "",
                       "imap-worker \u2500 email ingestion",
                       "    IMAP polling \u2192 ticket creation",
                       "",
                       "devops-worker \u2500 Azure DevOps sync",
                       "    Work item polling + AI workflow",
                       "",
                       "issue-resolver \u2500 auto code fixes",
                       "    BullMQ worker + Claude + git",
                       "",
                       "youtube-scheduler \u2500 stream setup",
                       "    Google Drive + YouTube API",
                       "",
                       "control-panel \u2500 Angular SPA",
                       "    Admin UI behind Caddy",
                       "",
                       "youtube-manager \u2500 YouTube admin",
                       "    Angular SPA bundled via Caddy",
                       "",
                       "status-monitor \u2500 health alerts",
                       "    Polls system status, sends alerts",
                       "",
                       "finance-app \u2500 personal finance",
                       "    Angular SPA at /finance/",
                       "",
                       "meal-planner \u2500 meal planning",
                       "    Angular SPA at /meals/",
                   ], accent=ACCENT_GREEN, title_size=16, item_size=12)

# mcp-servers/ + infra
add_component_card(slide, col3_x, top_y, card_w, Inches(2.0),
                   "mcp-servers/", [
                       "database \u2500 MCP database server",
                       "    Express + MCP SDK",
                       "    7 SQL tools for client DBs",
                       "    Runs in Azure App Service",
                   ], accent=ACCENT_TEAL, title_size=16, item_size=12)

add_component_card(slide, col3_x, Inches(4.8), card_w, Inches(1.9),
                   "Infrastructure", [
                       "docker-compose.yml \u2500 all services",
                       "Caddyfile \u2500 reverse proxy config",
                       ".github/workflows/ \u2500 CI/CD",
                       "tsconfig.base.json \u2500 shared TS config",
                   ], accent=ACCENT_ORANGE, title_size=16, item_size=12)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Data Model (Prisma Schema)
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_subtitle(slide, "Data Model", "PostgreSQL with Prisma ORM  \u00b7  packages/db/prisma/schema.prisma")

card_w = Inches(2.7)
card_h = Inches(2.0)
row1_y = Inches(2.5)
row2_y = Inches(4.8)

# Row 1: Core entities
add_component_card(slide, Inches(0.3), row1_y, card_w, card_h,
                   "Client", [
                       "id, name, shortCode",
                       "domainMappings[]",
                       "has: contacts, systems,",
                       "  tickets, codeRepos,",
                       "  integrations",
                   ], accent=ACCENT_BLUE, title_size=16, item_size=12)

add_component_card(slide, Inches(3.3), row1_y, card_w, card_h,
                   "Ticket", [
                       "clientId, systemId, subject",
                       "status: OPEN \u2192 CLOSED",
                       "priority: LOW \u2192 CRITICAL",
                       "source: EMAIL | DEVOPS | ...",
                       "category: DB_PERF | BUG | ...",
                       "has: events[], artifacts[]",
                   ], accent=ACCENT_BLUE, title_size=16, item_size=12)

add_component_card(slide, Inches(6.3), row1_y, card_w, card_h,
                   "TicketEvent", [
                       "ticketId, eventType, content",
                       "actor, metadata (JSON)",
                       "emailMessageId (dedup)",
                       "emailHash (SHA-256 dedup)",
                       "17 event types: COMMENT,",
                       "  AI_ANALYSIS, CODE_CHANGE...",
                   ], accent=ACCENT_BLUE, title_size=16, item_size=12)

add_component_card(slide, Inches(9.3), row1_y, Inches(3.5), card_h,
                   "System (Client DB)", [
                       "host, port, dbEngine",
                       "authMethod: SQL | Windows | AAD",
                       "environment: PROD | STAGING | ...",
                       "Connection pool config",
                       "has: tickets[], findings[],",
                       "  queryAuditLogs[]",
                   ], accent=ACCENT_TEAL, title_size=16, item_size=12)

# Row 2: Supporting entities
add_component_card(slide, Inches(0.3), row2_y, card_w, Inches(1.8),
                   "CodeRepo + IssueJob", [
                       "repoUrl, defaultBranch",
                       "branchPrefix (default: claude)",
                       "IssueJob: PENDING \u2192 COMPLETED",
                       "commitSha, filesChanged",
                   ], accent=ACCENT_GREEN, title_size=16, item_size=12)

add_component_card(slide, Inches(3.3), row2_y, card_w, Inches(1.8),
                   "DevOpsSyncState", [
                       "workItemId, revision",
                       "workflowState: idle \u2192 completed",
                       "planJson (structured steps)",
                       "Tracks conversation state",
                   ], accent=ACCENT_GREEN, title_size=16, item_size=12)

add_component_card(slide, Inches(6.3), row2_y, card_w, Inches(1.8),
                   "Finding + Playbook", [
                       "DB health findings",
                       "severity: LOW \u2192 CRITICAL",
                       "status: OPEN \u2192 RESOLVED",
                       "Playbook: remediation guide",
                   ], accent=ACCENT_PURPLE, title_size=16, item_size=12)

add_component_card(slide, Inches(9.3), row2_y, Inches(3.5), Inches(1.8),
                   "PromptOverride + Keyword", [
                       "Editable AI prompts from UI",
                       "scope: APP_WIDE | CLIENT",
                       "position: PREPEND | APPEND",
                       "Keywords: template placeholders",
                   ], accent=ACCENT_PURPLE, title_size=16, item_size=12)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 6 — copilot-api (Fastify)
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_subtitle(slide, "copilot-api  \u2014  REST API Gateway",
                   "services/copilot-api  \u00b7  Fastify  \u00b7  Port 3000  \u00b7  The central hub")

# Left side: description
add_bullet_list(slide, Inches(0.8), Inches(2.6), Inches(5.5), Inches(4.5), [
    "\u2022  Fastify HTTP server with JWT + API key authentication",
    "\u2022  CRUD endpoints for tickets, clients, systems, contacts, repos",
    "\u2022  BullMQ job producer \u2014 enqueues work for async workers",
    "\u2022  Proxies to MCP database server for SQL operations",
    "\u2022  Serves as the API backend for the Angular control panel",
    "\u2022  Health check at GET /api/health",
    "",
    "Key routes:",
    "  /api/tickets       \u2500  Ticket CRUD + event append",
    "  /api/clients       \u2500  Client management",
    "  /api/systems       \u2500  Database system registry",
    "  /api/repos         \u2500  Code repository CRUD",
    "  /api/issue-jobs    \u2500  Trigger issue resolution jobs",
    "  /api/auth          \u2500  Login, token refresh",
    "  /api/prompt-overrides  \u2500  AI prompt customisation",
], font_size=15, color=LIGHT_GRAY, spacing=Pt(5))

# Right side: connection diagram
card_x = Inches(7.2)
add_component_card(slide, card_x, Inches(2.5), Inches(5.5), Inches(1.2),
                   "Dependencies", [
                       "PostgreSQL (Prisma) \u00b7 Redis (BullMQ) \u00b7 Ollama \u00b7 Claude API",
                       "MCP Database Server (HTTP proxy to client SQL Servers)",
                   ], accent=ACCENT_BLUE, title_size=16, item_size=13)

add_component_card(slide, card_x, Inches(4.0), Inches(5.5), Inches(1.2),
                   "Authentication", [
                       "JWT tokens (control panel users) + API key header (service-to-service)",
                       "User roles: ADMIN, OPERATOR, CLIENT",
                   ], accent=ACCENT_GREEN, title_size=16, item_size=13)

add_component_card(slide, card_x, Inches(5.5), Inches(5.5), Inches(1.6),
                   "Docker", [
                       "Image: ghcr.io/siir/bronco/copilot-api",
                       "Depends on: postgres (healthy), redis (healthy)",
                       "Binds: 127.0.0.1:3000, QNAP mount for artifacts",
                       "Deployed to Hugo VM via deploy-hugo.yml workflow",
                   ], accent=ACCENT_ORANGE, title_size=16, item_size=13)

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 7 — imap-worker
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_subtitle(slide, "imap-worker  \u2014  Email Ingestion",
                   "services/imap-worker  \u00b7  IMAP polling  \u00b7  BullMQ consumer  \u00b7  Automatic ticket creation")

add_bullet_list(slide, Inches(0.8), Inches(2.6), Inches(5.5), Inches(4.5), [
    "How it works:",
    "",
    "1. Polls IMAP mailbox every 60s for new emails",
    "2. Encodes raw email as base64, pushes to BullMQ queue",
    "3. Processor parses email (mailparser)",
    "4. Idempotency: checks Message-ID + SHA-256 hash",
    "5. Threading logic:",
    "     a. In-Reply-To / References header matching",
    "     b. Subject line fuzzy match (7-day window)",
    "6. Creates new ticket or appends to existing thread",
    "7. Auto-provisions CLIENT user for known contacts",
    "8. Enqueues new tickets for AI analysis (triage, categorize)",
    "",
    "Key detail: Domain-based routing",
    "  \u2192  Maps sender email domain to Client.domainMappings[]",
    "  \u2192  Falls back to _unknown client if no match",
], font_size=15, color=LIGHT_GRAY, spacing=Pt(4))

# Right side flow
add_component_card(slide, Inches(7.0), Inches(2.5), Inches(5.8), Inches(4.5),
                   "Email Processing Pipeline", [
                       "",
                       "  \u2709  IMAP Server",
                       "    \u2502",
                       "    \u251c\u2500 Poll (every 60s)",
                       "    \u2502",
                       "    \u25bc",
                       "  BullMQ: email-processing queue",
                       "    \u2502",
                       "    \u251c\u2500 Parse email (simpleParser)",
                       "    \u251c\u2500 Dedup check (Message-ID + SHA-256)",
                       "    \u251c\u2500 Match sender \u2192 Contact \u2192 Client",
                       "    \u251c\u2500 Thread into existing ticket OR create new",
                       "    \u2502",
                       "    \u25bc",
                       "  PostgreSQL: Ticket + EMAIL_INBOUND event",
                       "    \u2502",
                       "    \u25bc",
                       "  BullMQ: analyze-ticket queue",
                       "    \u2502",
                       "    \u25bc",
                       "  AI: Triage + Categorize (Ollama)",
                   ], accent=ACCENT_GREEN, title_size=16, item_size=12)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 8 — devops-worker
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_subtitle(slide, "devops-worker  \u2014  Azure DevOps Integration",
                   "services/devops-worker  \u00b7  Work item sync  \u00b7  Conversational AI workflow engine")

add_bullet_list(slide, Inches(0.8), Inches(2.6), Inches(5.0), Inches(4.5), [
    "Polling & Sync:",
    "  \u2022  Polls Azure DevOps REST API (WIQL) every 120s",
    "  \u2022  Incremental sync (tracks revision per work item)",
    "  \u2022  Creates/updates tickets for each work item",
    "  \u2022  Syncs comments bidirectionally",
    "",
    "Conversational AI Workflow:",
    "  \u2022  Triggered for items assigned to configured user",
    "  \u2022  Full state machine with 7 states",
    "  \u2022  AI posts questions as DevOps comments",
    "  \u2022  User responds in DevOps \u2192 AI processes answers",
    "  \u2022  AI generates structured execution plan",
    "  \u2022  User approves/rejects via DevOps comments",
    "  \u2022  On approval: executes plan, posts results",
    "  \u2022  Intent classification: LLM with regex fallback",
], font_size=14, color=LIGHT_GRAY, spacing=Pt(4))

# Right side: state machine
add_component_card(slide, Inches(6.2), Inches(2.5), Inches(6.5), Inches(4.5),
                   "Workflow State Machine", [
                       "",
                       "  idle",
                       "    \u2502  New actionable work item detected",
                       "    \u25bc",
                       "  analyzing          \u2500  AI reads the work item & asks questions",
                       "    \u2502",
                       "    \u25bc",
                       "  questioning         \u2500  Waiting for user answers in DevOps",
                       "    \u2502                    (max 10 rounds, then force plan)",
                       "    \u25bc",
                       "  planning            \u2500  AI generates structured JSON plan",
                       "    \u2502",
                       "    \u25bc",
                       "  awaiting_approval   \u2500  Plan posted; waiting for LGTM/reject",
                       "    \u2502  \u250c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518",
                       "    \u2502  \u2502  Rejected \u2192 back to questioning",
                       "    \u25bc  \u2502  Clarification \u2192 answer & stay",
                       "  executing           \u2500  AI runs the approved plan",
                       "    \u2502",
                       "    \u25bc",
                       "  completed           \u2500  Results posted, ticket resolved",
                   ], accent=ACCENT_BLUE, title_size=16, item_size=12)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 9 — issue-resolver
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_subtitle(slide, "issue-resolver  \u2014  Automated Code Fixes",
                   "services/issue-resolver  \u00b7  BullMQ worker  \u00b7  Claude API  \u00b7  Git operations")

add_bullet_list(slide, Inches(0.8), Inches(2.6), Inches(5.5), Inches(4.5), [
    "How it works:",
    "",
    "1. Job triggered via POST /api/issue-jobs (ticketId + repoId)",
    "2. Worker picks up job from BullMQ queue",
    "3. Clones or pulls the target repository",
    "4. Gathers file tree + source context (up to 200KB of text files)",
    "5. Sends to Claude with issue description + full codebase context",
    "6. Claude returns JSON: {summary, changes: [{path, action, content}]}",
    "7. Applies changes locally with path traversal protection",
    "8. Commits to {branchPrefix}/{sanitized-slug} branch",
    "9. Pushes to remote (never to protected branches)",
    "10. Creates CODE_CHANGE ticket event with commit SHA",
    "",
    "Branch Safety (3 layers):",
    "  \u2022  API: rejects empty/protected branchPrefix",
    "  \u2022  Git: refuses to operate on main/master/develop/release",
    "  \u2022  Format: branch must contain '/' separator",
], font_size=14, color=LIGHT_GRAY, spacing=Pt(4))

# Right: job states
add_component_card(slide, Inches(7.0), Inches(2.5), Inches(5.8), Inches(2.5),
                   "IssueJob Status Flow", [
                       "PENDING \u2192 CLONING \u2192 ANALYZING \u2192 APPLYING \u2192 PUSHING \u2192 COMPLETED",
                       "",
                       "On failure at any stage: status = FAILED, error message saved",
                       "Each stage updates the job record in real time",
                   ], accent=ACCENT_GREEN, title_size=16, item_size=13)

add_component_card(slide, Inches(7.0), Inches(5.3), Inches(5.8), Inches(1.8),
                   "Claude Integration", [
                       "Model: claude-sonnet-4-5 (configurable)",
                       "Temperature: 0  (deterministic code generation)",
                       "Max tokens: 16,384  (large file output support)",
                       "Response: validated with Zod schema",
                   ], accent=ACCENT_PURPLE, title_size=16, item_size=13)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 10 — YouTube Scheduler
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_subtitle(slide, "youtube-scheduler  \u2014  Worship Stream Automation",
                   "services/youtube-scheduler  \u00b7  Google Drive API  \u00b7  YouTube Data API v3  \u00b7  Sharp")

add_bullet_list(slide, Inches(0.8), Inches(2.6), Inches(5.5), Inches(4.5), [
    "Pipeline (runs Saturday for Sunday service):",
    "",
    "1. Find latest .docx bulletin in Google Drive folder",
    "2. Convert to PDF via Drive API (copy as Google Doc \u2192 export PDF)",
    "3. Upload PDF back to Drive, share as \"anyone with link\"",
    "4. Extract theme paragraph from bulletin text (pages 2-3)",
    "5. Generate 1280\u00d7720 YouTube thumbnail:",
    "     \u2022  Composites theme text onto template image (Sharp library)",
    "     \u2022  Configurable font, colour, position",
    "6. Schedule YouTube live broadcast at configured day/time",
    "     \u2022  Creates broadcast + live stream + binds them",
    "     \u2022  Returns RTMP stream key for the encoder",
    "7. Upload custom thumbnail to the broadcast",
    "",
    "Idempotency: Persists last-scheduled date to JSON state file",
    "Polling: Checks every N hours if it's time to schedule",
], font_size=14, color=LIGHT_GRAY, spacing=Pt(4))

add_component_card(slide, Inches(7.0), Inches(2.5), Inches(5.8), Inches(4.5),
                   "Module Breakdown", [
                       "index.ts",
                       "  \u2502  Orchestrator, scheduling logic, config (Zod)",
                       "  \u2502  Timezone-aware date computation",
                       "  \u2502",
                       "bulletin-reader.ts",
                       "  \u2502  Drive API: find .docx, copy \u2192 Google Doc \u2192 PDF export",
                       "  \u2502  Upload PDF, set sharing permissions",
                       "  \u2502",
                       "theme-extractor.ts",
                       "  \u2502  Regex-based extraction from bulletin text",
                       "  \u2502  Targets pages 2-3, before order of service",
                       "  \u2502",
                       "thumbnail-generator.ts",
                       "  \u2502  Sharp: composite text SVG onto template image",
                       "  \u2502",
                       "youtube-broadcast.ts",
                       "    YouTube Data API v3: insert broadcast, stream, bind",
                       "    Upload custom thumbnail, return stream key + URL",
                   ], accent=ACCENT_ORANGE, title_size=16, item_size=12)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 11 — MCP Database Server
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_subtitle(slide, "MCP Database Server  \u2014  Client SQL Access",
                   "mcp-servers/database  \u00b7  Express + MCP SDK  \u00b7  Azure App Service  \u00b7  7 SQL tools")

add_bullet_list(slide, Inches(0.8), Inches(2.6), Inches(5.5), Inches(4.5), [
    "Purpose:",
    "  Provides AI-friendly SQL access to client database systems",
    "  Uses the Model Context Protocol (MCP) for tool registration",
    "",
    "Architecture:",
    "  \u2022  Express HTTP server (port 3100)",
    "  \u2022  Reads connection configs from SYSTEMS_CONFIG_PATH JSON file",
    "  \u2022  No dependency on the control plane PostgreSQL database",
    "  \u2022  Pool manager: factory pattern for mssql/tedious connections",
    "  \u2022  Credentials stored plaintext in JSON (Azure handles secrets)",
    "  \u2022  Query validator: SQL keyword blocklist (write protection)",
    "  \u2022  Audit logger: structured JSON to stdout (Pino)",
    "",
    "Extensibility:",
    "  \u2022  Add new DB engine: update DbEngine enum, add buildXxxConfig()",
    "  \u2022  Supported: MSSQL, Azure SQL MI, PostgreSQL, MySQL",
    "  \u2022  Connection pool per system, lazy initialisation",
], font_size=14, color=LIGHT_GRAY, spacing=Pt(4))

# Right: tool list
add_component_card(slide, Inches(7.0), Inches(2.5), Inches(5.8), Inches(4.5),
                   "MCP Tools (7)", [
                       "list_systems",
                       "  List all active client database systems",
                       "",
                       "run_query",
                       "  Execute read-only SQL (SELECT/WITH only, max 10K rows)",
                       "",
                       "inspect_schema",
                       "  List tables or inspect columns/constraints",
                       "",
                       "list_indexes",
                       "  Index details + optional usage statistics",
                       "",
                       "get_blocking_tree",
                       "  Current blocking chains with SQL text + wait info",
                       "",
                       "get_wait_stats",
                       "  Top N wait statistics (filters benign waits)",
                       "",
                       "get_database_health",
                       "  DB sizes, backup status, VLFs, CPU, memory, I/O latency",
                   ], accent=ACCENT_TEAL, title_size=16, item_size=12)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 12 — AI Routing
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_subtitle(slide, "AI Provider  \u2014  Dual-LLM Routing",
                   "packages/ai-provider  \u00b7  AIRouter class  \u00b7  Ollama (local) vs Claude (API)")

# Left: Local LLM
add_component_card(slide, Inches(0.5), Inches(2.5), Inches(5.8), Inches(4.5),
                   "Local LLM  \u2014  Ollama (Mac mini)", [
                       "Fast, cost-free inference for lightweight tasks",
                       "Model: llama3.1:8b (configurable)",
                       "Runs on siiriaplex Mac mini via Tailscale",
                       "",
                       "Tasks routed to local:",
                       "  TRIAGE            Set priority, extract entities",
                       "  CATEGORIZE        Classify into TicketCategory",
                       "  SUMMARIZE         Summarize email threads",
                       "  DRAFT_EMAIL       Generate draft responses",
                       "  EXTRACT_FACTS     Structured data from text",
                       "  SUMMARIZE_TICKET  Ticket summary generation",
                       "  SUGGEST_NEXT_STEPS  Workflow suggestions",
                       "  CLASSIFY_INTENT   Approval / rejection detection",
                       "  SUMMARIZE_LOGS    Summarize log entries",
                       "  ANALYZE_WORK_ITEM DevOps issue analysis",
                       "  DRAFT_COMMENT     Compose DevOps comments",
                       "  GENERATE_PLAN     Structured execution plans",
                       "  GENERATE_TITLE    Generate ticket titles",
                       "  CLASSIFY_EMAIL    Classify email intent/type",
                   ], accent=ACCENT_GREEN, title_size=16, item_size=12)

# Right: Claude API
add_component_card(slide, Inches(6.8), Inches(2.5), Inches(5.8), Inches(4.5),
                   "Claude API  \u2014  Deep Reasoning", [
                       "Heavy analysis requiring strong reasoning",
                       "Model: claude-sonnet-4-5 (configurable per service)",
                       "",
                       "Tasks routed to Claude:",
                       "  ANALYZE_QUERY       Query plan analysis",
                       "  GENERATE_SQL        SQL generation",
                       "  REVIEW_CODE         Code review",
                       "  DEEP_ANALYSIS       General deep analysis",
                       "  BUG_ANALYSIS        Cross-stack investigation",
                       "  ARCHITECTURE_REVIEW Architecture decisions",
                       "  SCHEMA_REVIEW       Schema change review",
                       "  FEATURE_ANALYSIS    Feature breakdown",
                       "  RESOLVE_ISSUE       Code generation for fixes",
                       "  CHANGE_CODEBASE_SMALL  Small code changes",
                       "  CHANGE_CODEBASE_LARGE  Large code changes",
                       "",
                       "Routing logic:",
                       "  if LOCAL_TASKS.has(taskType) \u2192 Ollama",
                       "  else \u2192 Claude (throws if key not configured)",
                   ], accent=ACCENT_PURPLE, title_size=16, item_size=12)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Infrastructure & Deployment
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_subtitle(slide, "Infrastructure & Deployment",
                   "Hugo VM  \u00b7  Azure App Service  \u00b7  Mac mini  \u00b7  Docker Compose  \u00b7  GitHub Actions")

# Hugo VM
add_component_card(slide, Inches(0.3), Inches(2.5), Inches(4.0), Inches(2.5),
                   "Hugo VM  (Ubuntu 24.04 LTS on ESXi NUC)", [
                       "Runs all core services via Docker Compose:",
                       "  \u2022 copilot-api (Fastify, port 3000)",
                       "  \u2022 imap-worker",
                       "  \u2022 devops-worker",
                       "  \u2022 issue-resolver",
                       "  \u2022 youtube-scheduler",
                       "  \u2022 PostgreSQL 16, Redis 7",
                       "  \u2022 Caddy (reverse proxy, Tailscale TLS)",
                   ], accent=ACCENT_BLUE, title_size=15, item_size=12)

# Azure
add_component_card(slide, Inches(4.6), Inches(2.5), Inches(4.0), Inches(2.5),
                   "Azure", [
                       "MCP Database Server",
                       "  \u2022 Azure App Service (ZIP deploy via publish profile)",
                       "  \u2022 Same vnet as client SQL MIs",
                       "  \u2022 Reads config from SYSTEMS_CONFIG_PATH JSON file",
                       "",
                       "Client SQL Managed Instances",
                       "  \u2022 SQL auth credentials",
                       "  \u2022 Connected via MCP server",
                   ], accent=ACCENT_TEAL, title_size=15, item_size=12)

# Mac mini
add_component_card(slide, Inches(8.9), Inches(2.5), Inches(4.0), Inches(2.5),
                   "Mac mini  (siiriaplex)", [
                       "Runs Ollama for local LLM inference",
                       "  \u2022 llama3.1:8b model",
                       "  \u2022 Accessible via Tailscale",
                       "  \u2022 Default: http://100.87.188.66:11434",
                       "",
                       "Cost-free AI for triage/classify tasks",
                       "Keeps latency low for high-volume work",
                   ], accent=ACCENT_PURPLE, title_size=15, item_size=12)

# CI/CD
add_component_card(slide, Inches(0.3), Inches(5.3), Inches(4.0), Inches(1.8),
                   "CI  \u2014  ci.yml", [
                       "Triggers: pushes + PRs to master",
                       "Auto-syncs pnpm-lock.yaml on PRs",
                       "Runs: typecheck + build (--frozen-lockfile)",
                   ], accent=ACCENT_ORANGE, title_size=15, item_size=12)

add_component_card(slide, Inches(4.6), Inches(5.3), Inches(4.0), Inches(1.8),
                   "Deploy Hugo  \u2014  deploy-hugo.yml", [
                       "Builds images \u2192 pushes to GHCR",
                       "SSHs to Hugo VM via Tailscale",
                       "docker compose pull + up -d",
                   ], accent=ACCENT_ORANGE, title_size=15, item_size=12)

add_component_card(slide, Inches(8.9), Inches(5.3), Inches(4.0), Inches(1.8),
                   "Deploy MCP  \u2014  deploy-mcp.yml", [
                       "Builds MCP database server",
                       "ZIP deploy to Azure App Service",
                       "Uses publish profile for auth",
                   ], accent=ACCENT_ORANGE, title_size=15, item_size=12)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Ticket Lifecycle
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_subtitle(slide, "Ticket Lifecycle  \u2014  End to End",
                   "From ingestion through AI triage to resolution")

add_component_card(slide, Inches(0.3), Inches(2.5), Inches(12.7), Inches(4.5),
                   "How a Ticket Flows Through the System", [
                       "",
                       "  \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500",
                       "  INGESTION                    AI TRIAGE                   WORK                        RESOLUTION",
                       "  \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500",
                       "",
                       "  Email lands in mailbox       Ollama: TRIAGE task         Operator reviews ticket     Manual resolution",
                       "    \u2502                            \u2502  Set priority             via control panel            \u2502  Operator works the issue",
                       "    \u25bc                            \u2502  Extract key entities         \u2502                          \u2502",
                       "  imap-worker picks it up      \u2502                              \u25bc                         OR",
                       "    \u2502  Parse + dedup            Ollama: CATEGORIZE task      Assign to system             \u2502",
                       "    \u2502  Thread or create ticket   \u2502  DATABASE_PERF?              \u2502                          \u25bc",
                       "    \u25bc                            \u2502  BUG_FIX?                    \u25bc                        Auto resolution",
                       "                                 \u2502  FEATURE_REQUEST?          Claude: DEEP_ANALYSIS         \u2502  issue-resolver",
                       "  Azure DevOps work item       \u2502  CODE_REVIEW?                \u2502  or BUG_ANALYSIS           \u2502  Claude generates code",
                       "    \u2502                            \u2502  ARCHITECTURE?               \u2502  or REVIEW_CODE            \u2502  Commits + pushes branch",
                       "    \u25bc                            \u25bc                              \u25bc                            \u25bc",
                       "  devops-worker syncs it       Category + Priority set     AI recommendations          CODE_CHANGE event",
                       "    \u2502  AI conversation                                       posted as ticket events      PR ready for review",
                       "    \u2502  Plan \u2192 Approve \u2192 Execute",
                       "",
                       "  Manual / API / Scheduled     Ollama: SUMMARIZE           Ticket status:              Ticket \u2192 RESOLVED",
                       "    \u2502  POST /api/tickets         Generate ticket summary     OPEN \u2192 IN_PROGRESS \u2192 RESOLVED \u2192 CLOSED",
                   ], accent=ACCENT_BLUE, title_size=16, item_size=11)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Key Conventions & Patterns
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_subtitle(slide, "Key Conventions & Patterns",
                   "What every developer needs to know")

add_component_card(slide, Inches(0.3), Inches(2.5), Inches(6.2), Inches(2.0),
                   "TypeScript Patterns", [
                       "ESM throughout \u2014 use .js extensions in relative imports",
                       "const object + type pattern for enums (no TS enums)",
                       "    export const Foo = { A: 'A', B: 'B' } as const;",
                       "    export type Foo = (typeof Foo)[keyof typeof Foo];",
                       "Zod for all config validation via loadConfig(schema)",
                       "Use z.output<typeof schema> (not z.infer) with .default()",
                   ], accent=ACCENT_BLUE, title_size=16, item_size=12)

add_component_card(slide, Inches(6.8), Inches(2.5), Inches(6.2), Inches(2.0),
                   "Logging & Observability", [
                       "Pino logging via createLogger(name) \u2192 stderr",
                       "AppLogger: structured logs written to Postgres app_logs table",
                       "AI token usage logging: provider, model, tokens, duration",
                       "Query audit: every SQL query logged with hash, caller, duration",
                   ], accent=ACCENT_GREEN, title_size=16, item_size=12)

add_component_card(slide, Inches(0.3), Inches(4.8), Inches(6.2), Inches(2.4),
                   "Database & Prisma", [
                       "Prisma ORM for control plane PostgreSQL",
                       "Prisma enum values must match shared-types exactly",
                       "All models use @map for snake_case DB column names",
                       "Soft-delete pattern: isActive flag on major entities",
                       "Dedup: unique constraints + hash-based idempotency",
                       "MCP server: independent (local JSON config, no control plane DB dependency)",
                   ], accent=ACCENT_PURPLE, title_size=16, item_size=12)

add_component_card(slide, Inches(6.8), Inches(4.8), Inches(6.2), Inches(2.4),
                   "Development Workflow", [
                       "pnpm install  \u2192  pnpm build  \u2192  pnpm typecheck",
                       "Lockfile discipline: always update pnpm-lock.yaml",
                       "CI: --frozen-lockfile (will fail on out-of-sync lockfile)",
                       "Node >= 20, pnpm >= 9",
                       "Branch safety: never push to main/master/develop/release",
                       "Issue resolver branches: {prefix}/{slug} format enforced",
                   ], accent=ACCENT_ORANGE, title_size=16, item_size=12)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Summary
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_subtitle(slide, "Summary", None, title_size=44)

summary_items = [
    "\u2022  Monorepo (pnpm workspaces): 4 shared packages, 7 services, 1 MCP server",
    "",
    "\u2022  copilot-api \u2014 central Fastify REST API, the hub that everything connects to",
    "\u2022  imap-worker \u2014 email \u2192 tickets with threading, dedup, and AI triage",
    "\u2022  devops-worker \u2014 Azure DevOps sync with full conversational AI workflow",
    "\u2022  issue-resolver \u2014 Claude-powered automated code fixes pushed to feature branches",
    "\u2022  youtube-scheduler \u2014 automated worship live stream scheduling via Google APIs",
    "\u2022  MCP database server \u2014 AI-friendly SQL access to client databases (Azure)",
    "",
    "\u2022  Dual AI: Ollama (fast/free local) for triage, Claude (API) for deep analysis",
    "\u2022  PostgreSQL (Prisma) for control plane, Redis + BullMQ for job queues",
    "\u2022  Docker Compose on Hugo VM, MCP in Azure App Service, Ollama on Mac mini",
    "\u2022  CI/CD via GitHub Actions: typecheck, build, deploy to GHCR + Azure",
]

add_bullet_list(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(5.0),
                summary_items, font_size=17, color=LIGHT_GRAY, spacing=Pt(5))


# ═════════════════════════════════════════════════════════════════════════
# Save
# ═════════════════════════════════════════════════════════════════════════
output_path = "docs/bronco-architecture.pptx"

import os
os.makedirs("docs", exist_ok=True)

prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"  Slides: {len(prs.slides)}")
